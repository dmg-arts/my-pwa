"""Builds the TOP-Feedback introduction deck.

    python3 -m venv .venv && .venv/bin/pip install python-pptx
    .venv/bin/python tools/docs/build-deck.py ./shots docs/TOP-Feedback-Introduction.pptx

Two rules keep this file openable everywhere, both learned the hard way:

  1. **Only fonts that exist on Windows, macOS and Google Slides.** Calibri ships
     with every Office install; Courier New is on all three. A deck that asks for
     a mac-only face renders with substituted metrics on the machine it is
     actually presented from, and every carefully placed box shifts.

  2. **Never let text exceed its box.** PowerPoint does not clip an overflowing
     text box — it spills the text over whatever is beneath it. Every block of
     text here goes through `fit()`, which shrinks the type until the estimate
     fits with room to spare, and `audit()` at the end re-checks the built file
     and refuses to finish if anything still overflows.
"""

import struct
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Inches, Pt

SHOTS = Path(sys.argv[1])
OUT = Path(sys.argv[2])

# --- palette, taken from the app's own stylesheet ---------------------------
NAVY = RGBColor(0x1C, 0x4F, 0x8B)
INK = RGBColor(0x14, 0x18, 0x1F)
MUTED = RGBColor(0x55, 0x60, 0x6E)
FAINT = RGBColor(0x8A, 0x93, 0x9F)
TINT = RGBColor(0xEC, 0xF2, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RULE = RGBColor(0xD3, 0xDC, 0xE8)
OK = RGBColor(0x1C, 0x6B, 0x3F)
WARN = RGBColor(0x8A, 0x5A, 0x05)
DANGER = RGBColor(0xA0, 0x27, 0x24)

SANS = "Calibri"
MONO = "Courier New"

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.9)                      # page margin
CONTENT_W = W - 2 * M

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]

EMU_IN = 914400


# ------------------------------------------------------------------ metrics


def est_height(text, size_pt, width_in, line_spacing=1.22):
    """Conservative height estimate for wrapped text, in inches.

    Calibri averages about 0.46 em per character for mixed-case prose. 0.54 is
    used instead so the estimate errs toward "needs more room" — the cost of
    being wrong that way is a slightly small font, and the cost of being wrong
    the other way is text lying on top of other text.
    """
    if not text:
        return 0.0
    chars_per_line = max(1, int(width_in / (size_pt * 0.54 / 72)))
    lines = 0
    for para in text.split("\n"):
        lines += max(1, -(-len(para) // chars_per_line))
    return lines * size_pt * line_spacing / 72


def fit(text, size_pt, width_in, height_in, floor_pt=11):
    """Largest size at or below `size_pt` whose estimate fits the box."""
    size = size_pt
    while size > floor_pt and est_height(text, size, width_in) > height_in:
        size -= 1
    return size


# ------------------------------------------------------------------ drawing


def slide():
    return prs.slides.add_slide(BLANK)


def shape(s, kind, x, y, w, h, fill_rgb=None, line_rgb=None, radius=None):
    sp = s.shapes.add_shape(kind, int(x), int(y), int(w), int(h))
    if radius is not None and kind == MSO_SHAPE.ROUNDED_RECTANGLE:
        sp.adjustments[0] = radius
    if fill_rgb is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill_rgb
    if line_rgb is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line_rgb
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    sp.text_frame.text = ""          # keep a valid, empty text body
    return sp


def rect(s, x, y, w, h, fill_rgb=None, line_rgb=None, radius=None):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE
    return shape(s, kind, x, y, w, h, fill_rgb, line_rgb, radius)


def background(s, rgb):
    rect(s, 0, 0, W, H, rgb)


def block(s, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, shrink=True):
    """Places paragraphs in a box, shrinking the type until they fit.

    `paras` is a list of (text, size_pt, bold, colour, space_after_pt).
    """
    width_in = w / EMU_IN
    height_in = h / EMU_IN

    if shrink:
        # Scale every paragraph by one common factor so the hierarchy survives.
        total = sum(est_height(t, sz, width_in) + after / 72 for t, sz, _, _, after in paras)
        step = 0
        while total > height_in and step < 12:
            step += 1
            total = sum(
                est_height(t, max(10, sz - step), width_in) + after / 72
                for t, sz, _, _, after in paras
            )
        paras = [(t, max(10, sz - step), b, c, after) for t, sz, b, c, after in paras]

    box = s.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    for i, (body, size, bold, colour, after) in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(after)
        p.line_spacing = 1.22
        r = p.add_run()
        r.text = body
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = colour
        r.font.name = SANS
    return box


def points(s, x, y, w, h, items, size=17, gap=13):
    """A list where a leading '-' marks a supporting line under the point."""
    paras = []
    for item in items:
        sub = item.startswith("-")
        body = item[1:].strip() if sub else item
        paras.append((("   " + body) if sub else body,
                      size - 3 if sub else size,
                      not sub,
                      MUTED if sub else INK,
                      gap - 5 if sub else gap))
    return block(s, x, y, w, h, paras)


def eyebrow(s, label, colour=NAVY):
    block(s, M, Inches(0.55), Inches(9), Inches(0.32),
          [(label.upper(), 12, True, colour, 0)], shrink=False)


def title(s, text, sub=None, colour=INK):
    size = fit(text, 34, CONTENT_W / EMU_IN, 1.05)
    block(s, M, Inches(0.98), CONTENT_W, Inches(1.05),
          [(text, size, True, colour, 0)], shrink=False)
    if sub:
        sz = fit(sub, 16, 10.8, 0.62)
        block(s, M, Inches(1.95), Inches(10.9), Inches(0.62),
              [(sub, sz, False, MUTED, 0)], shrink=False)


def png_size(path):
    with open(path, "rb") as fh:
        return struct.unpack(">II", fh.read(24)[16:24])


def picture(s, name, x, y, box_w, box_h, frame=True):
    """Scales an image to fit its box without distortion, and centres it."""
    path = SHOTS / name
    pw, ph = png_size(path)
    scale = min(box_w / pw, box_h / ph)
    w, h = int(pw * scale), int(ph * scale)
    left, top = int(x + (box_w - w) / 2), int(y + (box_h - h) / 2)
    if frame:
        pad = 9525
        rect(s, left - pad, top - pad, w + 2 * pad, h + 2 * pad, WHITE, RULE)
    s.shapes.add_picture(str(path), left, top, w, h)


def page_no(s, n):
    block(s, Inches(11.9), Inches(6.92), Inches(0.55), Inches(0.3),
          [(str(n), 11, False, FAINT, 0)], align=PP_ALIGN.RIGHT, shrink=False)


def notes(s, body):
    s.notes_slide.notes_text_frame.text = body


def card(s, x, y, w, h, heading, body, accent=NAVY):
    rect(s, x, y, w, h, WHITE, RULE, radius=0.04)
    rect(s, x, y, w, Inches(0.055), accent)
    pad = Inches(0.34)
    block(s, x + pad, y + Inches(0.42), w - 2 * pad, h - Inches(0.72),
          [(heading, 19, True, INK, 10), (body, 14, False, MUTED, 0)])


# ------------------------------------------------------------------- slides

def s01_title():
    s = slide()
    background(s, NAVY)
    rect(s, M, Inches(1.35), Inches(0.9), Inches(0.9), WHITE, radius=0.2)
    block(s, M, Inches(1.47), Inches(0.9), Inches(0.66),
          [("✓", 38, True, NAVY, 0)], align=PP_ALIGN.CENTER, shrink=False)

    block(s, M, Inches(2.85), Inches(11), Inches(1.25),
          [("TOP-Feedback", 58, True, WHITE, 0)], shrink=False)
    block(s, M, Inches(4.2), Inches(9.8), Inches(1.1),
          [("The student and instructor feedback cycle,", 22, False, WHITE, 5),
           ("owned entirely by your detachment.", 22, False, WHITE, 0)], shrink=False)
    rect(s, M, Inches(5.72), Inches(1.15), Emu(28575), WHITE)
    block(s, M, Inches(6.0), Inches(10), Inches(0.45),
          [("An introduction — what it is, what it does, and how it is used", 15, False, WHITE, 0)],
          shrink=False)
    notes(s, "A primer, not a training session. Three things to land: there is no vendor and no "
             "server; the detachment owns the data; and it turns feedback into something an "
             "instructor can act on rather than a stack of paper.")


def s02_what():
    s = slide()
    eyebrow(s, "What it is")
    title(s, "A feedback app that stores nothing on anyone else's computer")

    block(s, M, Inches(2.12), Inches(11.4), Inches(1.0),
          [("Cadets answer a short form on their phone. Instructors see the results analysed, "
            "not just listed. Every record is written to a Google Drive folder your detachment "
            "already owns.", 18, False, MUTED, 0)])

    w = Inches(3.55)
    gap = Inches(0.42)
    x = M
    for head, body in [
        ("Your Google account",
         "No vendor holds your data, because no vendor is involved. If you stop using the app, "
         "every response is still sitting in your Drive."),
        ("Any device",
         "Runs in a browser and installs to a phone home screen. Works with no signal and sends "
         "when the connection returns."),
        ("Nothing to buy",
         "No licences, no subscription, no server to maintain. Setup is a Drive folder and about "
         "45 minutes, once."),
    ]:
        card(s, x, Inches(3.3), w, Inches(2.75), head, body)
        x += w + gap
    page_no(s, 2)
    notes(s, "These three answer the questions always asked first: where does the data go, what do "
             "I need to run it, and what does it cost.")


def s03_problem():
    s = slide()
    eyebrow(s, "Why it exists")
    title(s, "Feedback nobody can act on")

    points(s, M, Inches(2.3), Inches(5.9), Inches(3.9), [
        "Paper forms get collected, stacked, and never counted.",
        "Survey tools put cadet feedback on a company's servers.",
        "An average hides the very thing you needed to see.",
        "A disclosure of hazing can sit unread in a pile for weeks.",
    ], size=18, gap=22)

    rect(s, Inches(7.25), Inches(2.2), Inches(5.2), Inches(3.5), TINT, radius=0.04)
    block(s, Inches(7.7), Inches(2.62), Inches(4.3), Inches(2.7),
          [("What it replaces", 16, True, NAVY, 12),
           ("One app that issues the form, collects the responses, runs the statistics, "
            "reads the written answers, and flags anything a person needs to see today.",
            15, False, INK, 12),
           ("On a phone. Offline. At no cost.", 15, True, INK, 0)])
    page_no(s, 3)
    notes(s, "Frame the problem before the product. The last bullet is the one a commander cares "
             "about most.")


def s04_data():
    s = slide()
    eyebrow(s, "Where the data lives")
    title(s, "Your folder is the database",
          "Each record is a plain file you can open in Drive — nothing is locked in a format only this app can read.")

    rect(s, M, Inches(2.75), Inches(6.0), Inches(3.5), TINT, radius=0.03)
    tree = ("TOP-Feedback/\n"
            "   config/      detachment profile and settings\n"
            "   users/       accounts: cadets, instructors, admins\n"
            "   forms/       feedback form definitions\n"
            "   requests/    feedback issued, each with an ID\n"
            "   responses/   submitted feedback\n"
            "   receipts/    who submitted, kept apart from\n"
            "                what they said\n"
            "   reports/     exported reports")
    box = s.shapes.add_textbox(M + Inches(0.35), Inches(3.0), Inches(5.4), Inches(3.0))
    tf = box.text_frame
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(tree.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(3)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(11)
        r.font.name = MONO
        r.font.color.rgb = NAVY if line.rstrip().endswith("/") else MUTED

    points(s, Inches(7.4), Inches(2.85), Inches(5.0), Inches(3.4), [
        "Stop using the app and you still have every response.",
        "-Readable in Drive, backed up by Google, recoverable from the trash for 30 days.",
        "Who can read the feedback is a Drive sharing setting.",
        "-Not a setting inside this app. That is the access control that actually matters.",
    ], size=16, gap=16)
    page_no(s, 4)
    notes(s, "This is the procurement answer: no lock-in, no data processor, no contract. Stress "
             "that folder sharing is the real access control.")


def s05_roles():
    s = slide()
    eyebrow(s, "Who uses it")
    title(s, "One app, three doors")
    picture(s, "home.png", M, Inches(2.2), Inches(6.5), Inches(4.15))
    points(s, Inches(7.9), Inches(2.45), Inches(4.6), Inches(3.9), [
        "Student",
        "-Signs in, sees only the feedback assigned to them, submits it once.",
        "Instructor Portal",
        "-Creates feedback, reads responses, runs the analysis.",
        "Database Administration",
        "-Creates accounts and resets passwords.",
    ], size=16, gap=13)
    page_no(s, 5)
    notes(s, "Everyone signs in, cadets included. That is what makes one-submission-per-student "
             "a fact rather than an honour system.")


def s06_student_list():
    s = slide()
    eyebrow(s, "For cadets")
    title(s, "Only what is assigned to you",
          "Filtered by school year, semester, class and due date. Anything already submitted drops off the list.")
    picture(s, "student-list.png", M, Inches(2.75), Inches(11.5), Inches(3.85))
    page_no(s, 6)
    notes(s, "A cadet sees their own list, not everything the detachment ever issued. Scheduled "
             "forms stay hidden until they open; overdue ones close on their own.")


def s07_scale():
    s = slide()
    eyebrow(s, "Answering")
    title(s, "Cadets choose a word. The maths runs on the number.")
    picture(s, "student-form.png", M, Inches(2.35), Inches(7.1), Inches(4.3))
    points(s, Inches(8.4), Inches(2.5), Inches(4.1), Inches(4.0), [
        "Nine ratings, every one named",
        "-Detrimental, Significant, Unfavorable, Minor, Neutral, Slight, Favorable, Major, Outstanding.",
        "No numbers on screen",
        "-A visible number invites people to average it in their head while answering.",
        "Written answers capped at 250 words",
        "-With a live counter as they type.",
    ], size=15, gap=11)
    page_no(s, 7)
    notes(s, "Neutral sits at 5, the true centre of a 1-9 scale. Instructors see the word and the "
             "number together; cadets only ever see the word.")


def s08_phone():
    s = slide()
    background(s, TINT)
    eyebrow(s, "On a phone")
    title(s, "Built for the walk out of the classroom")
    points(s, M, Inches(2.55), Inches(6.4), Inches(3.9), [
        "Installs to the home screen",
        "-Opens like an ordinary app. No address for anyone to remember.",
        "Works with no signal",
        "-A submission made offline is saved on the device, shown as pending, and sent "
        "automatically when the connection returns.",
        "Nothing is lost to a dropped connection",
    ], size=17, gap=15)
    picture(s, "student-mobile.png", Inches(8.5), Inches(0.85), Inches(3.9), Inches(5.9), frame=True)
    page_no(s, 8)
    notes(s, "Classroom wifi is unreliable — that is the whole reason for the offline queue. "
             "Reading existing feedback still needs a connection.")


def s09_portal():
    s = slide()
    eyebrow(s, "For instructors")
    title(s, "Two things to do, both one click away")
    picture(s, "instructor-portal.png", M, Inches(2.4), Inches(11.5), Inches(4.15))
    page_no(s, 9)
    notes(s, "Create Feedback and Feedback Response and Analysis are the headline actions. The "
             "tabs underneath hold forms, students and database tools.")


def s10_create():
    s = slide()
    eyebrow(s, "Creating feedback")
    title(s, "A standardized form, in about two minutes")
    picture(s, "form-creator.png", M, Inches(2.35), Inches(6.8), Inches(4.3))
    points(s, Inches(8.1), Inches(2.45), Inches(4.4), Inches(4.15), [
        "Name the class or event",
        "Set the AS level, term and dates",
        "-Scheduled forms stay hidden until they open.",
        "Choose who receives it",
        "-Everyone at that level, or named cadets.",
        "Add questions — three minimum",
        "-A rating or a written answer. No maximum.",
        "Every form gets an ID",
        "-Such as FB-2026-0001, to quote in conversation.",
    ], size=14, gap=10)
    page_no(s, 10)
    notes(s, "The form is deliberately narrow. That constraint is exactly what makes results "
             "comparable between events, cohorts and terms.")


def s11_ratings():
    s = slide()
    eyebrow(s, "Analysis — ratings")
    title(s, "Everything an average on its own would hide")
    picture(s, "analysis-question.png", M, Inches(2.35), Inches(7.2), Inches(4.3))
    points(s, Inches(8.5), Inches(2.45), Inches(4.0), Inches(4.15), [
        "Mean, median, range and spread",
        "The distribution across every rating",
        "-Shape you can see rather than infer.",
        "An agreement score from 0 to 1",
        "-Did they actually agree, or just average out?",
        "Outliers, found robustly",
        "-Measured against the median, so one extreme rating cannot hide itself.",
    ], size=14, gap=11)
    page_no(s, 11)
    notes(s, "Point at the distribution in the screenshot. It is bimodal — the mean of 5.75 "
             "describes nobody who was actually in the room.")


def s12_split():
    s = slide()
    background(s, NAVY)
    block(s, M, Inches(1.45), Inches(11.4), Inches(2.5),
          [("An average of 5 can mean everyone shrugged —", 33, True, WHITE, 10),
           ("or that half the flight thought it outstanding", 33, True, WHITE, 10),
           ("and half thought it harmful.", 33, True, WHITE, 0)], shrink=False)
    rect(s, M, Inches(4.35), Inches(1.3), Emu(28575), WHITE)
    block(s, M, Inches(4.75), Inches(10.8), Inches(1.9),
          [("Same number. Opposite situations. Only one of them needs acting on.",
            19, False, WHITE, 14),
           ("TOP-Feedback finds the split, says so plainly, and tells you the average describes "
            "nobody — so you read both groups instead of the middle.", 17, False, WHITE, 0)])
    notes(s, "This is the strongest single argument for the analysis over a spreadsheet. A "
             "spreadsheet gives you the mean and stops there.")


def s13_written():
    s = slide()
    eyebrow(s, "Analysis — written answers")
    title(s, "Read the ones that matter first")
    points(s, M, Inches(2.35), Inches(5.9), Inches(4.2), [
        "Sentiment, ranked most negative first",
        "-Handles negation and emphasis, so 'not helpful' and 'extremely helpful' both score correctly.",
        "Every answer is shown",
        "-Answers the word list cannot read are labelled, never hidden.",
        "It runs on the device",
        "-No cloud service sees your cadets' words. That is why it is a word list and not a model, "
        "and the accuracy limits are stated on screen.",
    ], size=16, gap=13)
    picture(s, "wordcloud.png", Inches(7.05), Inches(2.35), Inches(5.4), Inches(3.8))
    block(s, Inches(7.05), Inches(6.3), Inches(5.4), Inches(0.5),
          [("Word cloud sized by how many people used a word, not how often it appears.",
            12, False, FAINT, 0)], shrink=False)
    page_no(s, 13)
    notes(s, "Be straight about the limits in the room: it reads words, not meaning. It decides "
             "what you read first; it does not read for you.")


def s14_safety():
    s = slide()
    eyebrow(s, "Safety screen", DANGER)
    title(s, "Nothing important sits unread")
    picture(s, "safety.png", M, Inches(2.35), Inches(7.4), Inches(4.3))
    points(s, Inches(8.7), Inches(2.45), Inches(3.8), Inches(4.15), [
        "Every written answer is screened",
        "-Hazing, sexual harassment, discrimination, violence, self-harm, substances, integrity.",
        "Matches are highlighted in context",
        "A prompt, never a finding",
        "-It cannot tell 'we discussed hazing prevention' from 'I was hazed'. A person reads it "
        "and decides.",
    ], size=14, gap=11)
    page_no(s, 14)
    notes(s, "Say the limits out loud: it matches words, not meaning, and a clear screen is not "
             "proof that nothing was reported. Follow detachment reporting procedures either way.")


def s15_anonymity():
    s = slide()
    eyebrow(s, "Anonymity")
    title(s, "Anonymous, and still able to chase people")

    w = Inches(3.55)
    gap = Inches(0.42)
    x = M
    for head, body, accent in [
        ("How it works",
         "The cadet signs in, so the app knows who they are. It writes a receipt in a separate "
         "folder, then drops the name. The response itself carries no identity.",
         NAVY),
        ("What that buys",
         "You can say 'two cadets still owe feedback' about responses nobody can attribute — and "
         "nobody can submit twice.",
         OK),
        ("The honest limit",
         "With very few responses a single answer can be identified by elimination. So anonymous "
         "results stay hidden until three people have responded.",
         WARN),
    ]:
        card(s, x, Inches(2.4), w, Inches(3.05), head, body, accent)
        x += w + gap

    block(s, M, Inches(5.85), Inches(11.4), Inches(0.85),
          [("Safety flags are the deliberate exception: a disclosure cannot wait for a third "
            "response, so the app tells you one exists and warns that opening it may identify "
            "the author.", 15, False, INK, 0)])
    page_no(s, 15)
    notes(s, "Expect this question. The receipt-versus-response separation is the whole design, "
             "and the three-response threshold is the honest limit of the maths.")


def s16_admin():
    s = slide()
    eyebrow(s, "Running it")
    title(s, "Accounts your detachment controls",
          "Add cadets one at a time or import a whole class by CSV. Passwords are generated and handed to you once.")
    picture(s, "admin.png", M, Inches(2.8), Inches(11.5), Inches(3.75))
    page_no(s, 16)
    notes(s, "Mention that the credentials list can only be read once, because only hashes are "
             "stored. Download it, hand it out, delete the file.")


def s17_start():
    s = slide()
    background(s, NAVY)
    block(s, M, Inches(1.1), Inches(10), Inches(0.95),
          [("Getting started", 42, True, WHITE, 0)], shrink=False)

    y = Inches(2.45)
    for n, head, body in [
        ("1", "Create the folder",
         "A Google account the detachment owns, and a folder named TOP-Feedback."),
        ("2", "Register the app",
         "One free Google Cloud project, so the app can reach that Drive. About ten minutes."),
        ("3", "Run the wizard",
         "Paste the folder link, create your administrator account, add cadets."),
    ]:
        circ = shape(s, MSO_SHAPE.OVAL, M, y, Inches(0.6), Inches(0.6), WHITE)
        block(s, M, y + Inches(0.09), Inches(0.6), Inches(0.42),
              [(n, 19, True, NAVY, 0)], align=PP_ALIGN.CENTER, shrink=False)
        block(s, M + Inches(0.95), y - Inches(0.02), Inches(9.6), Inches(1.05),
              [(head, 21, True, WHITE, 4), (body, 15, False, WHITE, 0)])
        y += Inches(1.25)

    rect(s, M, Inches(6.35), Inches(11.4), Emu(19050), WHITE)
    block(s, M, Inches(6.6), Inches(11.4), Inches(0.45),
          [("Step by step, with screenshots: TOP-Feedback Setup Guide (PDF).",
            15, False, WHITE, 0)], shrink=False)
    notes(s, "Close here. Anyone who will actually perform the install should be given the setup "
             "guide PDF rather than this deck.")


for build in [s01_title, s02_what, s03_problem, s04_data, s05_roles, s06_student_list,
              s07_scale, s08_phone, s09_portal, s10_create, s11_ratings, s12_split,
              s13_written, s14_safety, s15_anonymity, s16_admin, s17_start]:
    build()

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))


# -------------------------------------------------------------------- audit

def audit(path):
    """Re-open the built file and check no text can overflow its box."""
    import re
    import xml.etree.ElementTree as ET
    import zipfile

    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main",
          "p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
    z = zipfile.ZipFile(path)
    problems = []
    fonts = set()

    for name in z.namelist():
        if not re.match(r"ppt/slides/slide\d+\.xml$", name):
            continue
        root = ET.fromstring(z.read(name))
        fonts.update(re.findall(r'typeface="([^"]+)"', z.read(name).decode()))
        for sp in root.iter(f'{{{ns["p"]}}}sp'):
            ext = sp.find(".//a:ext", ns)
            tx = sp.find(".//p:txBody", ns)
            if ext is None or tx is None:
                continue
            cx, cy = int(ext.get("cx")), int(ext.get("cy"))
            if cx <= 0 or cy <= 0:
                continue
            need = 0.0
            for para in tx.findall("a:p", ns):
                runs = para.findall("a:r", ns)
                if not runs:
                    continue
                body = "".join((r.find("a:t", ns).text or "") for r in runs)
                sizes = [int(r.find("a:rPr", ns).get("sz", "1800"))
                         for r in runs if r.find("a:rPr", ns) is not None]
                size = max(sizes) / 100 if sizes else 18
                spc = para.find("a:pPr/a:spcAft/a:spcPts", ns)
                need += est_height(body, size, cx / EMU_IN)
                need += (int(spc.get("val")) / 100 / 72) if spc is not None else 0
            if need > cy / EMU_IN:
                problems.append(f"{name.split('/')[-1]}: needs {need:.2f}in, box {cy / EMU_IN:.2f}in")
    return problems, sorted(fonts)


issues, used_fonts = audit(OUT)
print(f"{len(prs.slides._sldIdLst)} slides -> {OUT}")
print(f"fonts used: {', '.join(used_fonts)}")
if issues:
    print(f"OVERFLOW in {len(issues)} box(es):")
    for i in issues:
        print("  " + i)
    sys.exit(1)
print("layout audit: no text overflows its box")
